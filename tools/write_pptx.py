"""
WAT tool: Generate a PowerPoint (.pptx) file from a JSON content spec.

Input JSON schema (--input or --input-file):
{
  "output_path": "path/to/output.pptx",
  "slide_width_inches": 13.33,   // optional, default 13.33 (widescreen 16:9)
  "slide_height_inches": 7.5,    // optional, default 7.5
  "slides": [
    {
      "layout": "title",          // "title" | "content" | "two_col" | "blank"
      "title": "Slide title",
      "subtitle": "Subtitle text",          // for layout=title
      "body": "Body text or bullet points", // plain string → one paragraph
      "bullets": ["Point 1", "Point 2"],    // list → bulleted list (layout=content)
      "left_col":  ["Left bullet 1", ...],  // layout=two_col
      "right_col": ["Right bullet 1", ...], // layout=two_col
      "notes": "Speaker notes text"
    }
  ]
}
"""

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv

load_dotenv()

# Slide layout indices for the default theme
_LAYOUT = {
    "title":   0,   # Title Slide
    "content": 1,   # Title and Content
    "two_col": 3,   # Two Content
    "blank":   6,   # Blank
}


def _add_textbox(slide, text, left, top, width, height, font_size=18, bold=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    return txBox


def _fill_text_frame(tf, bullets: list):
    for i, item in enumerate(bullets):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = item
        para.level = 0


def build_presentation(spec: dict) -> str:
    output_path = spec.get("output_path")
    if not output_path:
        raise ValueError("'output_path' is required in the input spec")

    prs = Presentation()
    prs.slide_width = Inches(spec.get("slide_width_inches", 13.33))
    prs.slide_height = Inches(spec.get("slide_height_inches", 7.5))

    slide_layouts = prs.slide_layouts

    for slide_spec in spec.get("slides", []):
        layout_name = slide_spec.get("layout", "content")
        layout_idx = _LAYOUT.get(layout_name, 1)
        slide_layout = slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

        # Title (idx 0 is always title)
        if slide_spec.get("title") and 0 in placeholders:
            placeholders[0].text = slide_spec["title"]

        if layout_name == "title":
            # idx 1 = subtitle on title slide
            if slide_spec.get("subtitle") and 1 in placeholders:
                placeholders[1].text = slide_spec["subtitle"]
            elif slide_spec.get("body") and 1 in placeholders:
                placeholders[1].text = slide_spec["body"]

        elif layout_name == "content":
            # idx 1 = content placeholder
            bullets = slide_spec.get("bullets") or (
                [slide_spec["body"]] if slide_spec.get("body") else []
            )
            if bullets and 1 in placeholders:
                _fill_text_frame(placeholders[1].text_frame, bullets)

        elif layout_name == "two_col":
            left = slide_spec.get("left_col", [])
            right = slide_spec.get("right_col", [])
            # idx 1 = left content, idx 2 = right content
            if left and 1 in placeholders:
                _fill_text_frame(placeholders[1].text_frame, left)
            if right and 2 in placeholders:
                _fill_text_frame(placeholders[2].text_frame, right)

        elif layout_name == "blank":
            if slide_spec.get("body"):
                _add_textbox(slide, slide_spec["body"],
                             left=0.5, top=1.0, width=12.0, height=5.5)

        # Speaker notes
        if slide_spec.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_spec["notes"]

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Write a PowerPoint .pptx file from a JSON spec")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--input", help="JSON string spec")
    group.add_argument("--input-file", help="Path to JSON spec file")
    parser.add_argument("--output-file", help="Optional path to write result JSON")
    args = parser.parse_args()

    try:
        if args.input_file:
            with open(args.input_file) as f:
                spec = json.load(f)
        else:
            spec = json.loads(args.input)

        saved_path = build_presentation(spec)
        result = {"status": "ok", "output_path": saved_path}
    except Exception as e:
        print(json.dumps({"status": "error", "message": str(e)}), file=sys.stderr)
        sys.exit(1)

    output = json.dumps(result, indent=2)
    print(output)

    if args.output_file:
        os.makedirs(os.path.dirname(args.output_file) or ".", exist_ok=True)
        with open(args.output_file, "w") as f:
            f.write(output)


if __name__ == "__main__":
    main()
